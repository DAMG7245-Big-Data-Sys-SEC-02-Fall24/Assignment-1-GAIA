import json
import pandas as pd
import matplotlib.pyplot as plt
import os
from dotenv import load_dotenv
from typing import List, Union
import re
from PIL import Image
import io
from openai import OpenAI
import requests
import zipfile
import PyPDF2
import docx
import pptx
import wave
import mutagen
import base64
import streamlit as st
load_dotenv()

client = OpenAI(
    base_url=os.getenv("base_url") or st.secrets("base_url"),
    api_key=os.getenv("api_key") or st.secrets("api_key"))



class FileProcessor:
    def __init__(self):
        pass

    def search(self, query: str) -> str:
        """Performs a web search using DuckDuckGo."""
        url = f"https://api.duckduckgo.com/?q={query}&format=json"
        response = requests.get(url)
        results = response.json()
        return results.get('Abstract', 'No results found.')

    def read_image(self, file_path: str) -> str:
        """Reads an image file and returns its base64 encoded string."""
        try:
            with open(file_path, "rb") as image_file:
                return base64.b64encode(image_file.read()).decode('utf-8')
        except Exception as e:
            return f"Error reading image: {str(e)}"

    def read_excel(self, file_path: str) -> str:
        """Reads an Excel file and returns its content as a JSON string."""
        try:
            df = pd.read_excel(file_path)
            print(df)
            return df.to_json(orient='records')
        except Exception as e:
            return f"Error reading Excel file: {str(e)}"

    def read_csv(self, file_path: str) -> str:
        """Reads a CSV file and returns its content as a JSON string."""
        try:
            df = pd.read_csv(file_path)
            return df.to_json(orient='records')
        except Exception as e:
            return f"Error reading CSV file: {str(e)}"

    def read_zip(self, file_path: str) -> str:
        """Reads a ZIP file and returns a list of its contents."""
        try:
            with zipfile.ZipFile(file_path, 'r') as zip_ref:
                return json.dumps(zip_ref.namelist())
        except Exception as e:
            return f"Error reading ZIP file: {str(e)}"

    def read_pdf(self, file_path: str) -> str:
        """Reads a PDF file and returns its text content."""
        try:
            with open(file_path, 'rb') as file:
                reader = PyPDF2.PdfReader(file)
                text = ""
                for page in reader.pages:
                    page_text = page.extract_text()
                    if page_text:
                        text += page_text + "\n"
            return text
        except Exception as e:
            return f"Error reading PDF file: {str(e)}"

    def read_json(self, file_path: str) -> str:
        """Reads a JSON file and returns its content."""
        try:
            with open(file_path, 'r') as file:
                return json.dumps(json.load(file))
        except Exception as e:
            return f"Error reading JSON file: {str(e)}"

    def read_python(self, file_path: str) -> str:
        """Reads a Python file and returns its content."""
        try:
            with open(file_path, 'r') as file:
                return file.read()
        except Exception as e:
            return f"Error reading Python file: {str(e)}"

    def read_docx(self, file_path: str) -> str:
        """Reads a DOCX file and returns its text content."""
        try:
            doc = docx.Document(file_path)
            return "\n".join([paragraph.text for paragraph in doc.paragraphs])
        except Exception as e:
            return f"Error reading DOCX file: {str(e)}"

    def read_pptx(self, file_path: str) -> str:
        """Reads a PPTX file and returns its text content."""
        try:
            prs = pptx.Presentation(file_path)
            text = ""
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, 'text'):
                        text += shape.text + "\n"
            return text
        except Exception as e:
            return f"Error reading PPTX file: {str(e)}"

    def read_audio(self, file_path: str) -> str:
        """Reads an audio file, transcribes it using OpenAI's Whisper API,
        and returns its transcription along with metadata."""
        try:
            if file_path.endswith('.mp3') or file_path.endswith('.wav'):
                # Extract metadata using mutagen for MP3 or wave for WAV
                metadata = {}
                transcription = ""
                if file_path.endswith('.mp3'):
                    audio = mutagen.File(file_path)
                    metadata = {
                        "length": audio.info.length,
                        "bitrate": audio.info.bitrate
                    }
                elif file_path.endswith('.wav'):
                    with wave.open(file_path, 'rb') as wav:
                        metadata = {
                            "channels": wav.getnchannels(),
                            "sample_width": wav.getsampwidth(),
                            "framerate": wav.getframerate(),
                            "frames": wav.getnframes()
                        }

                # Transcribe the audio using OpenAI's Whisper API
                with open(file_path, "rb") as audio_file:
                    transcription_response = client.audio.transcribe(model="whisper-1",
                    file=audio_file)
                    transcription = transcription_response.get("text", "")

                metadata["transcription"] = transcription
                return json.dumps(metadata)
            else:
                return "Unsupported audio format for transcription."
        except Exception as e:
            return f"Error reading/transcribing audio file: {str(e)}"

    def read_pdb(self, file_path: str) -> str:
        """Reads a PDB file and returns its content."""
        try:
            with open(file_path, 'r') as file:
                return file.read()
        except Exception as e:
            return f"Error reading PDB file: {str(e)}"

    def read_jsonld(self, file_path: str) -> str:
        """Reads a JSON-LD file and returns its content."""
        try:
            with open(file_path, 'r') as file:
                return json.dumps(json.load(file))
        except Exception as e:
            return f"Error reading JSON-LD file: {str(e)}"

    def read_txt(self, file_path: str) -> str:
        """Reads a TXT file and returns its content."""
        try:
            with open(file_path, 'r', encoding='utf-8') as file:
                return file.read()
        except Exception as e:
            return f"Error reading TXT file: {str(e)}"


# Define the tools
file_processor = FileProcessor()
tools = {
    "Search": file_processor.search,
    "ReadImage": file_processor.read_image,
    "ReadExcel": file_processor.read_excel,
    "ReadCSV": file_processor.read_csv,
    "ReadZIP": file_processor.read_zip,
    "ReadPDF": file_processor.read_pdf,
    "ReadJSON": file_processor.read_json,
    "ReadPython": file_processor.read_python,
    "ReadDOCX": file_processor.read_docx,
    "ReadPPTX": file_processor.read_pptx,
    "ReadAudio": file_processor.read_audio,
    "ReadPDB": file_processor.read_pdb,
    "ReadJSONLD": file_processor.read_jsonld,
    "ReadTXT": file_processor.read_txt  # Added ReadTXT
}
